from pptx import Presentation

# Create a presentation object
presentation = Presentation()

# Define slide titles and contents
slides_content = [
    # Title Slide
    {
        "title": "Simple E-commerce System",
        "content": ["Presented by: Your Name", "Date: [Insert Date]"]
    },
    # Introduction Slide
    {
        "title": "Introduction",
        "content": [
            "E-commerce allows buying and selling online.",
            "Aim: To create a simple, user-friendly system.",
            "Focus: Customers, Products, and Transactions."
        ]
    },
    # Problem Statement Slide
    {
        "title": "Problem Statement",
        "content": [
            "Traditional shopping is time-consuming.",
            "Limited accessibility for remote customers.",
            "Need for a reliable online shopping solution."
        ]
    },
    # Scope Slide
    {
        "title": "Scope of the Project",
        "content": [
            "Supports product listing and searching.",
            "Facilitates order placement and tracking.",
            "Focuses on a basic user-friendly interface."
        ]
    },
    # Methodology Slide
    {
        "title": "Methodology",
        "content": [
            "Technologies: Python, Django, SQLite.",
            "Frontend: HTML, CSS, JavaScript.",
            "Backend: User Authentication, Product Management."
        ]
    },
    # System Design Slide
    {
        "title": "System Design",
        "content": [
            "High-Level Design:",
            "- User Module: Registration/Login",
            "- Product Module: Add/View Products",
            "- Order Module: Place/Track Orders"
        ]
    },
    # Implementation Slide
    {
        "title": "Implementation",
        "content": [
            "Frontend: Responsive web design.",
            "Backend: Secure database integration.",
            "Demo: Basic workflow for adding products and placing orders."
        ]
    },
    # Testing Slide
    {
        "title": "Testing",
        "content": [
            "Performed functional and unit testing.",
            "Example: Adding products, user login functionality.",
            "Outcome: Identified and resolved minor issues."
        ]
    },
    # Results Slide
    {
        "title": "Results and Analysis",
        "content": [
            "Fully functional simple e-commerce system.",
            "User-friendly and accessible interface.",
            "Potential for future scaling and enhancements."
        ]
    },
    # Challenges Slide
    {
        "title": "Challenges and Limitations",
        "content": [
            "Challenges: Limited time for extensive testing.",
            "Limitations: Basic features, lacks advanced filters.",
            "Future Plans: Adding payment gateway and detailed analytics."
        ]
    },
    # Conclusion Slide
    {
        "title": "Conclusion",
        "content": [
            "Successfully developed a basic e-commerce system.",
            "Achieved user-friendly design and functionality.",
            "Open for enhancements and feature expansion."
        ]
    },
    # Thank You Slide
    {
        "title": "Thank You",
        "content": ["Questions and Discussions Welcome!"]
    }
]

# Add slides to the presentation
for slide_info in slides_content:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content layout
    slide.shapes.title.text = slide_info["title"]
    content_box = slide.placeholders[1]
    content_box.text = "\n".join(slide_info["content"])

# Save the presentation
file_path = "/Users/anilthakur/Desktop/python/Simple_Ecommerce_System_Presentation.pptx"
presentation.save(file_path)
file_path